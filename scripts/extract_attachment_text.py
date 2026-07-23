#!/usr/bin/env python3
from __future__ import annotations

import argparse
import gzip
import hashlib
import json
import os
import re
import shutil
import subprocess
import tempfile
import time
from collections import Counter
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable, Iterator
from urllib.parse import unquote, urlsplit

import requests
from docx import Document as DocxDocument
from docx.document import Document as DocxDocumentType
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph
from odf import teletype
from odf.namespaces import TABLENS, TEXTNS
from odf.opendocument import load as load_odt
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry


QUEUE_PATH = Path(
    "public/rag/attachment-extraction-queue-final.jsonl"
)
MANIFEST_PATH = Path(
    "public/rag/attachment-manifest.jsonl"
)
UNREGISTERED_PATH = Path(
    "public/rag/unregistered-attachment-links.jsonl"
)

INDEX_PATH = Path(
    "public/rag/attachment-text-index.jsonl"
)
STATE_PATH = Path(
    "public/rag/attachment-extraction-state.jsonl"
)
ERRORS_PATH = Path(
    "public/rag/attachment-extraction-errors.jsonl"
)
CHANGES_PATH = Path(
    "public/changes/attachment-text-changes.jsonl"
)
STATUS_PATH = Path(
    "public/attachment-extraction-status.json"
)
BLOB_DIRECTORY = Path(
    "public/rag/attachment-text-blobs"
)

DEFAULT_MAX_FILE_SIZE_MB = 30
DEFAULT_TIMEOUT_SECONDS = 90
DEFAULT_WORKERS = 1

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".doc",
    ".docx",
    ".odt",
    ".ppt",
    ".pptx",
    ".xls",
    ".xlsx",
}

OFFICE_ZIP_EXTENSIONS = {
    ".docx",
    ".odt",
    ".pptx",
    ".xlsx",
}

LEGACY_OFFICE_EXTENSIONS = {
    ".doc",
    ".ppt",
    ".xls",
}

SCHEMA_VERSION = 1


@dataclass
class ExtractedContent:
    text: str
    sections: list[dict[str, Any]]
    parser: str
    status: str
    warnings: list[str]
    metrics: dict[str, Any]


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat()


def load_jsonl(path: Path) -> list[dict[str, Any]]:
    if not path.exists():
        return []

    records: list[dict[str, Any]] = []

    with path.open("r", encoding="utf-8") as file:
        for line_number, raw_line in enumerate(
            file,
            start=1,
        ):
            line = raw_line.strip()

            if not line:
                continue

            try:
                record = json.loads(line)
            except json.JSONDecodeError as error:
                raise ValueError(
                    f"Nieprawidłowy JSON w {path}, "
                    f"linia {line_number}: {error}"
                ) from error

            if isinstance(record, dict):
                records.append(record)

    return records


def write_jsonl(
    path: Path,
    records: Iterable[dict[str, Any]],
) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)

    with path.open("w", encoding="utf-8") as file:
        for record in records:
            file.write(
                json.dumps(
                    record,
                    ensure_ascii=False,
                    separators=(",", ":"),
                )
            )
            file.write("\n")


def normalize_url(value: Any) -> str:
    return str(value or "").strip()


def file_extension(
    url: str,
    filename: str | None = None,
) -> str:
    if filename:
        suffix = Path(filename).suffix.lower()

        if suffix:
            return suffix

    path = unquote(urlsplit(url).path)
    return Path(path).suffix.lower()


def safe_filename(
    value: str | None,
    fallback: str,
) -> str:
    filename = str(value or "").strip()

    if not filename:
        filename = fallback

    filename = filename.replace("\\", "_")
    filename = filename.replace("/", "_")
    filename = re.sub(r"[\x00-\x1f\x7f]+", "_", filename)
    filename = filename.strip(" .")

    return filename or fallback


def normalize_text(value: Any) -> str:
    text = str(value or "")
    text = text.replace("\x00", "")
    text = text.replace("\r\n", "\n")
    text = text.replace("\r", "\n")

    cleaned_lines: list[str] = []
    previous_blank = False

    for raw_line in text.split("\n"):
        line = re.sub(r"[ \t]+", " ", raw_line).strip()

        if not line:
            if not previous_blank:
                cleaned_lines.append("")
            previous_blank = True
            continue

        cleaned_lines.append(line)
        previous_blank = False

    return "\n".join(cleaned_lines).strip()


def word_count(text: str) -> int:
    return len(
        re.findall(
            r"\b[\wÀ-ž'-]+\b",
            text,
            flags=re.UNICODE,
        )
    )


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()

    with path.open("rb") as file:
        for chunk in iter(
            lambda: file.read(1024 * 1024),
            b"",
        ):
            digest.update(chunk)

    return digest.hexdigest()


def queue_signature(item: dict[str, Any]) -> str:
    payload = {
        "effective_url": normalize_url(
            item.get("effective_url")
            or item.get("final_url")
            or item.get("url")
        ),
        "filename": item.get("filename"),
        "declared_filesize": item.get(
            "declared_filesize"
        ),
        "content_length": item.get(
            "content_length"
        ),
        "declared_mime_type": item.get(
            "declared_mime_type"
        ),
        "content_type": item.get(
            "content_type"
        ),
        "resolution": item.get(
            "resolution"
        ),
    }

    serialized = json.dumps(
        payload,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
    ).encode("utf-8")

    return sha256_bytes(serialized)


def relationship_key(item: dict[str, Any]) -> tuple[str, str, str]:
    return (
        str(item.get("document_id") or ""),
        str(item.get("document_url") or ""),
        str(item.get("anchor_text") or ""),
    )


def build_relationship_context() -> dict[str, dict[str, Any]]:
    context: dict[str, dict[str, Any]] = {}

    def ensure(url: str) -> dict[str, Any]:
        return context.setdefault(
            url,
            {
                "linked_from": [],
                "attachment_metadata": {},
            },
        )

    for item in load_jsonl(MANIFEST_PATH):
        urls = {
            normalize_url(item.get("url")),
            normalize_url(item.get("normalized_url")),
        }
        urls.discard("")

        metadata = {
            "attachment_registry_id": item.get("id"),
            "wordpress_id": item.get("wordpress_id"),
            "attachment_page_url": item.get(
                "attachment_page_url"
            ),
            "title": item.get("title"),
            "caption": item.get("caption"),
            "description": item.get("description"),
            "parent_wordpress_id": item.get(
                "parent_wordpress_id"
            ),
            "published_at": item.get("published_at"),
            "modified_at": item.get("modified_at"),
            "classification": item.get("classification"),
        }

        linked_from = [
            relation
            for relation in item.get("linked_from", [])
            if isinstance(relation, dict)
        ]

        for url in urls:
            target = ensure(url)
            target["attachment_metadata"] = {
                key: value
                for key, value in metadata.items()
                if value is not None
            }
            target["linked_from"].extend(linked_from)

    for item in load_jsonl(UNREGISTERED_PATH):
        url = normalize_url(item.get("url"))

        if not url:
            continue

        target = ensure(url)
        target["linked_from"].extend(
            relation
            for relation in item.get("linked_from", [])
            if isinstance(relation, dict)
        )

    for target in context.values():
        unique: dict[
            tuple[str, str, str],
            dict[str, Any],
        ] = {}

        for relation in target["linked_from"]:
            unique[relationship_key(relation)] = relation

        target["linked_from"] = sorted(
            unique.values(),
            key=lambda relation: (
                str(relation.get("document_source_id") or ""),
                str(relation.get("document_title") or ""),
                str(relation.get("document_url") or ""),
            ),
        )

    return context


def resolve_relationship_context(
    relationship_context: dict[str, dict[str, Any]],
    item: dict[str, Any],
    effective_url: str,
) -> dict[str, Any]:
    urls = {
        effective_url,
        normalize_url(item.get("url")),
        normalize_url(item.get("requested_url")),
        normalize_url(item.get("final_url")),
        normalize_url(item.get("replacement_url")),
    }
    urls.update(
        normalize_url(alias)
        for alias in item.get("url_aliases", [])
    )
    urls.discard("")

    linked_from: dict[
        tuple[str, str, str],
        dict[str, Any],
    ] = {}
    attachment_metadata: dict[str, Any] = {}

    for url in urls:
        context = relationship_context.get(url)

        if not context:
            continue

        if not attachment_metadata:
            attachment_metadata = dict(
                context.get("attachment_metadata") or {}
            )

        for relation in context.get("linked_from", []):
            linked_from[relationship_key(relation)] = relation

    relations = sorted(
        linked_from.values(),
        key=lambda relation: (
            str(relation.get("document_source_id") or ""),
            str(relation.get("document_title") or ""),
            str(relation.get("document_url") or ""),
        ),
    )

    return {
        "linked_from": relations,
        "linked_from_count": len(relations),
        "attachment_metadata": attachment_metadata,
    }


def create_session() -> requests.Session:
    session = requests.Session()

    retry = Retry(
        total=4,
        connect=4,
        read=3,
        status=4,
        backoff_factor=1.0,
        status_forcelist={
            429,
            500,
            502,
            503,
            504,
        },
        allowed_methods={"GET"},
        respect_retry_after_header=True,
    )

    adapter = HTTPAdapter(
        max_retries=retry,
        pool_connections=4,
        pool_maxsize=4,
    )

    session.mount("https://", adapter)
    session.mount("http://", adapter)

    session.headers.update(
        {
            "User-Agent": (
                "UEW-RAG-Document-Extractor/1.0"
            ),
            "Accept": "*/*",
        }
    )

    return session


def parse_content_length(
    headers: requests.structures.CaseInsensitiveDict,
) -> int | None:
    value = headers.get("Content-Length")

    if value and value.isdigit():
        return int(value)

    content_range = headers.get(
        "Content-Range",
        "",
    )

    match = re.search(
        r"/(\d+)$",
        content_range,
    )

    if match:
        return int(match.group(1))

    return None


def download_document(
    *,
    session: requests.Session,
    url: str,
    destination: Path,
    previous_state: dict[str, Any] | None,
    timeout_seconds: int,
    max_bytes: int,
    force: bool,
) -> dict[str, Any]:
    headers: dict[str, str] = {}

    if previous_state and not force:
        etag = str(
            previous_state.get("etag")
            or ""
        ).strip()
        last_modified = str(
            previous_state.get(
                "last_modified"
            )
            or ""
        ).strip()

        if etag:
            headers["If-None-Match"] = etag

        if last_modified:
            headers["If-Modified-Since"] = (
                last_modified
            )

    with session.get(
        url,
        headers=headers,
        stream=True,
        allow_redirects=True,
        timeout=(
            20,
            timeout_seconds,
        ),
    ) as response:
        if response.status_code == 304:
            return {
                "download_status": "not_modified",
                "http_status": 304,
                "final_url": response.url or url,
                "etag": (
                    response.headers.get("ETag")
                    or (
                        previous_state.get("etag")
                        if previous_state
                        else None
                    )
                ),
                "last_modified": (
                    response.headers.get(
                        "Last-Modified"
                    )
                    or (
                        previous_state.get(
                            "last_modified"
                        )
                        if previous_state
                        else None
                    )
                ),
                "content_length": (
                    previous_state.get(
                        "content_length"
                    )
                    if previous_state
                    else None
                ),
                "content_type": (
                    previous_state.get(
                        "content_type"
                    )
                    if previous_state
                    else None
                ),
                "redirect_count": len(
                    response.history
                ),
            }

        response.raise_for_status()

        content_length = parse_content_length(
            response.headers
        )

        if (
            content_length is not None
            and content_length > max_bytes
        ):
            raise ValueError(
                "Plik przekracza limit rozmiaru: "
                f"{content_length} > {max_bytes} B"
            )

        destination.parent.mkdir(
            parents=True,
            exist_ok=True,
        )

        downloaded_bytes = 0

        with destination.open("wb") as file:
            for chunk in response.iter_content(
                chunk_size=1024 * 1024,
            ):
                if not chunk:
                    continue

                downloaded_bytes += len(chunk)

                if downloaded_bytes > max_bytes:
                    raise ValueError(
                        "Plik przekracza limit rozmiaru "
                        f"{max_bytes} B podczas pobierania."
                    )

                file.write(chunk)

        return {
            "download_status": "downloaded",
            "http_status": response.status_code,
            "final_url": response.url or url,
            "etag": response.headers.get("ETag"),
            "last_modified": response.headers.get(
                "Last-Modified"
            ),
            "content_length": (
                content_length
                if content_length is not None
                else downloaded_bytes
            ),
            "content_type": response.headers.get(
                "Content-Type"
            ),
            "redirect_count": len(
                response.history
            ),
            "downloaded_bytes": downloaded_bytes,
        }


def section(
    *,
    section_id: str,
    kind: str,
    label: str,
    text: str,
    metadata: dict[str, Any] | None = None,
) -> dict[str, Any] | None:
    cleaned = normalize_text(text)

    if not cleaned:
        return None

    result = {
        "section_id": section_id,
        "kind": kind,
        "label": label,
        "text": cleaned,
        "char_count": len(cleaned),
        "word_count": word_count(cleaned),
    }

    if metadata:
        result["metadata"] = metadata

    return result


def extract_pdf(path: Path) -> ExtractedContent:
    reader = PdfReader(str(path))

    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception:
            pass

    sections: list[dict[str, Any]] = []
    warnings: list[str] = []
    failed_pages: list[int] = []

    for page_number, page in enumerate(
        reader.pages,
        start=1,
    ):
        try:
            page_text = page.extract_text() or ""
        except Exception as error:
            failed_pages.append(page_number)
            warnings.append(
                "Nie udało się odczytać strony "
                f"{page_number}: {error}"
            )
            continue

        item = section(
            section_id=f"page-{page_number}",
            kind="page",
            label=f"Strona {page_number}",
            text=page_text,
            metadata={
                "page_number": page_number,
            },
        )

        if item:
            sections.append(item)

    text = normalize_text(
        "\n\n".join(
            item["text"]
            for item in sections
        )
    )

    status = (
        "ok"
        if len(text) >= 20
        else "needs_ocr"
    )

    if status == "needs_ocr":
        warnings.append(
            "PDF nie zawiera wystarczającej "
            "warstwy tekstowej; może wymagać OCR."
        )

    return ExtractedContent(
        text=text,
        sections=sections,
        parser="pypdf",
        status=status,
        warnings=warnings,
        metrics={
            "page_count": len(reader.pages),
            "text_page_count": len(sections),
            "failed_page_count": len(
                failed_pages
            ),
            "failed_pages": failed_pages,
        },
    )


def iter_docx_blocks(
    document: DocxDocumentType,
) -> Iterator[
    DocxParagraph | DocxTable
]:
    body = document.element.body

    for child in body.iterchildren():
        tag = child.tag.rsplit("}", 1)[-1]

        if tag == "p":
            yield DocxParagraph(
                child,
                document,
            )
        elif tag == "tbl":
            yield DocxTable(
                child,
                document,
            )


def docx_table_text(table: DocxTable) -> str:
    rows: list[str] = []

    for row in table.rows:
        values = [
            normalize_text(cell.text)
            for cell in row.cells
        ]

        if any(values):
            rows.append(
                "\t".join(values)
            )

    return "\n".join(rows)


def extract_docx(path: Path) -> ExtractedContent:
    document = DocxDocument(str(path))
    sections: list[dict[str, Any]] = []
    warnings: list[str] = []

    current_label = "Treść dokumentu"
    current_parts: list[str] = []
    section_number = 1

    def flush_current() -> None:
        nonlocal current_parts
        nonlocal section_number

        item = section(
            section_id=(
                f"section-{section_number}"
            ),
            kind="document_section",
            label=current_label,
            text=(
                (current_label + "\n")
                if current_label != "Treść dokumentu"
                else ""
            )
            + "\n".join(current_parts),
        )

        if item:
            sections.append(item)
            section_number += 1

        current_parts = []

    for block in iter_docx_blocks(document):
        if isinstance(block, DocxParagraph):
            paragraph_text = normalize_text(
                block.text
            )

            if not paragraph_text:
                continue

            style_name = str(
                block.style.name
                if block.style is not None
                else ""
            ).lower()

            if style_name.startswith("heading"):
                flush_current()
                current_label = paragraph_text
                continue

            current_parts.append(paragraph_text)

        elif isinstance(block, DocxTable):
            table_text = docx_table_text(block)

            if table_text:
                current_parts.append(
                    table_text
                )

    flush_current()

    header_footer_parts: list[str] = []

    for section_index, doc_section in enumerate(
        document.sections,
        start=1,
    ):
        header_text = normalize_text(
            "\n".join(
                paragraph.text
                for paragraph
                in doc_section.header.paragraphs
            )
        )
        footer_text = normalize_text(
            "\n".join(
                paragraph.text
                for paragraph
                in doc_section.footer.paragraphs
            )
        )

        if header_text:
            header_footer_parts.append(
                f"Nagłówek sekcji "
                f"{section_index}\n{header_text}"
            )

        if footer_text:
            header_footer_parts.append(
                f"Stopka sekcji "
                f"{section_index}\n{footer_text}"
            )

    if header_footer_parts:
        item = section(
            section_id="headers-footers",
            kind="headers_footers",
            label="Nagłówki i stopki",
            text="\n\n".join(
                header_footer_parts
            ),
        )

        if item:
            sections.append(item)

    text = normalize_text(
        "\n\n".join(
            item["text"]
            for item in sections
        )
    )

    status = (
        "ok"
        if len(text) >= 5
        else "empty_text"
    )

    return ExtractedContent(
        text=text,
        sections=sections,
        parser="python-docx",
        status=status,
        warnings=warnings,
        metrics={
            "section_count": len(sections),
            "table_count": len(
                document.tables
            ),
        },
    )


def extract_shape_text(shape: Any) -> list[str]:
    values: list[str] = []

    if getattr(
        shape,
        "has_text_frame",
        False,
    ):
        text = normalize_text(
            shape.text
        )

        if text:
            values.append(text)

    if getattr(
        shape,
        "has_table",
        False,
    ):
        for row in shape.table.rows:
            cells = [
                normalize_text(cell.text)
                for cell in row.cells
            ]

            if any(cells):
                values.append(
                    "\t".join(cells)
                )

    if getattr(
        shape,
        "shape_type",
        None,
    ) == 6:
        for child_shape in shape.shapes:
            values.extend(
                extract_shape_text(
                    child_shape
                )
            )

    return values


def extract_pptx(path: Path) -> ExtractedContent:
    presentation = Presentation(str(path))
    sections: list[dict[str, Any]] = []
    warnings: list[str] = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        parts: list[str] = []

        for shape in slide.shapes:
            parts.extend(
                extract_shape_text(shape)
            )

        try:
            notes_slide = slide.notes_slide
            notes_text = normalize_text(
                "\n".join(
                    shape.text
                    for shape
                    in notes_slide.shapes
                    if getattr(
                        shape,
                        "has_text_frame",
                        False,
                    )
                )
            )

            if notes_text:
                parts.append(
                    "Notatki prelegenta:\n"
                    f"{notes_text}"
                )

        except Exception as error:
            warnings.append(
                "Nie udało się odczytać notatek "
                f"slajdu {slide_number}: {error}"
            )

        item = section(
            section_id=f"slide-{slide_number}",
            kind="slide",
            label=f"Slajd {slide_number}",
            text="\n".join(parts),
            metadata={
                "slide_number": slide_number,
            },
        )

        if item:
            sections.append(item)

    text = normalize_text(
        "\n\n".join(
            item["text"]
            for item in sections
        )
    )

    status = (
        "ok"
        if len(text) >= 5
        else "empty_text"
    )

    return ExtractedContent(
        text=text,
        sections=sections,
        parser="python-pptx",
        status=status,
        warnings=warnings,
        metrics={
            "slide_count": len(
                presentation.slides
            ),
            "text_slide_count": len(
                sections
            ),
        },
    )


def cell_to_text(
    formula_value: Any,
    cached_value: Any,
) -> str:
    if cached_value not in {
        None,
        "",
    }:
        return normalize_text(cached_value)

    if formula_value not in {
        None,
        "",
    }:
        return normalize_text(formula_value)

    return ""


def extract_xlsx(path: Path) -> ExtractedContent:
    workbook_formula = load_workbook(
        filename=str(path),
        read_only=True,
        data_only=False,
    )
    workbook_values = load_workbook(
        filename=str(path),
        read_only=True,
        data_only=True,
    )

    sections: list[dict[str, Any]] = []
    warnings: list[str] = []
    sheet_names = list(workbook_formula.sheetnames)

    try:
        for sheet_name in sheet_names:
            formula_sheet = workbook_formula[
                sheet_name
            ]
            value_sheet = workbook_values[
                sheet_name
            ]

            rows: list[str] = []
            non_empty_cells = 0

            for formula_row, value_row in zip(
                formula_sheet.iter_rows(),
                value_sheet.iter_rows(),
            ):
                values: list[str] = []

                for formula_cell, value_cell in zip(
                    formula_row,
                    value_row,
                ):
                    cell_text = cell_to_text(
                        formula_cell.value,
                        value_cell.value,
                    )
                    values.append(cell_text)

                    if cell_text:
                        non_empty_cells += 1

                while values and not values[-1]:
                    values.pop()

                if any(values):
                    rows.append(
                        "\t".join(values)
                    )

            item = section(
                section_id=(
                    f"sheet-{len(sections) + 1}"
                ),
                kind="worksheet",
                label=sheet_name,
                text="\n".join(rows),
                metadata={
                    "sheet_name": sheet_name,
                    "non_empty_cell_count": (
                        non_empty_cells
                    ),
                },
            )

            if item:
                sections.append(item)

    finally:
        workbook_formula.close()
        workbook_values.close()

    text = normalize_text(
        "\n\n".join(
            item["text"]
            for item in sections
        )
    )

    status = (
        "ok"
        if len(text) >= 5
        else "empty_text"
    )

    return ExtractedContent(
        text=text,
        sections=sections,
        parser="openpyxl",
        status=status,
        warnings=warnings,
        metrics={
            "sheet_count": len(sheet_names),
            "text_sheet_count": len(
                sections
            ),
        },
    )


def iter_odt_blocks(node: Any) -> Iterator[Any]:
    for child in getattr(node, "childNodes", []):
        qname = getattr(child, "qname", None)

        if qname in {
            (TEXTNS, "h"),
            (TEXTNS, "p"),
            (TABLENS, "table"),
        }:
            yield child

            if qname == (TABLENS, "table"):
                continue

        yield from iter_odt_blocks(child)


def odt_table_text(table_node: Any) -> str:
    rows: list[str] = []

    for row in getattr(table_node, "childNodes", []):
        if getattr(row, "qname", None) != (
            TABLENS,
            "table-row",
        ):
            continue

        values: list[str] = []

        for cell in getattr(row, "childNodes", []):
            if getattr(cell, "qname", None) not in {
                (TABLENS, "table-cell"),
                (TABLENS, "covered-table-cell"),
            }:
                continue

            value = normalize_text(
                teletype.extractText(cell)
            )
            values.append(value)

        while values and not values[-1]:
            values.pop()

        if any(values):
            rows.append("\t".join(values))

    return "\n".join(rows)


def extract_odt(path: Path) -> ExtractedContent:
    document = load_odt(str(path))
    sections: list[dict[str, Any]] = []
    warnings: list[str] = []

    current_label = "Treść dokumentu"
    current_parts: list[str] = []
    section_number = 1

    def flush_current() -> None:
        nonlocal current_parts
        nonlocal section_number

        item = section(
            section_id=(
                f"section-{section_number}"
            ),
            kind="document_section",
            label=current_label,
            text=(
                (current_label + "\n")
                if current_label != "Treść dokumentu"
                else ""
            )
            + "\n".join(current_parts),
        )

        if item:
            sections.append(item)
            section_number += 1

        current_parts = []

    for node in iter_odt_blocks(document.text):
        qname = getattr(node, "qname", None)

        if qname == (TEXTNS, "h"):
            heading = normalize_text(
                teletype.extractText(node)
            )

            if heading:
                flush_current()
                current_label = heading
            continue

        if qname == (TEXTNS, "p"):
            paragraph = normalize_text(
                teletype.extractText(node)
            )

            if paragraph:
                current_parts.append(paragraph)
            continue

        if qname == (TABLENS, "table"):
            table_text = odt_table_text(node)

            if table_text:
                current_parts.append(table_text)

    flush_current()

    text = normalize_text(
        "\n\n".join(
            item["text"]
            for item in sections
        )
    )

    status = (
        "ok"
        if len(text) >= 5
        else "empty_text"
    )

    return ExtractedContent(
        text=text,
        sections=sections,
        parser="odfpy",
        status=status,
        warnings=warnings,
        metrics={
            "section_count": len(sections),
        },
    )


def extract_legacy_doc(
    path: Path,
) -> ExtractedContent:
    antiword = shutil.which("antiword")

    if antiword:
        process = subprocess.run(
            [
                antiword,
                str(path),
            ],
            capture_output=True,
            check=False,
        )

        raw = (
            process.stdout.decode(
                "utf-8",
                errors="replace",
            )
        )

        if (
            process.returncode != 0
            and not raw.strip()
        ):
            error_text = process.stderr.decode(
                "utf-8",
                errors="replace",
            )
            raise RuntimeError(
                "antiword nie odczytał pliku: "
                f"{error_text}"
            )

        text = normalize_text(raw)
        item = section(
            section_id="document",
            kind="document",
            label="Treść dokumentu",
            text=text,
        )
        sections = [item] if item else []

        return ExtractedContent(
            text=text,
            sections=sections,
            parser="antiword",
            status=(
                "ok"
                if len(text) >= 5
                else "empty_text"
            ),
            warnings=[],
            metrics={
                "section_count": len(
                    sections
                ),
            },
        )

    libreoffice = (
        shutil.which("libreoffice")
        or shutil.which("soffice")
    )

    if libreoffice:
        with tempfile.TemporaryDirectory() as directory:
            output_directory = Path(directory)

            process = subprocess.run(
                [
                    libreoffice,
                    "--headless",
                    "--convert-to",
                    "txt:Text",
                    "--outdir",
                    str(output_directory),
                    str(path),
                ],
                capture_output=True,
                text=True,
                check=False,
            )

            output_path = output_directory / (
                path.stem + ".txt"
            )

            if (
                process.returncode != 0
                or not output_path.exists()
            ):
                raise RuntimeError(
                    "LibreOffice nie odczytał "
                    f"pliku: {process.stderr}"
                )

            text = normalize_text(
                output_path.read_text(
                    encoding="utf-8",
                    errors="replace",
                )
            )

        item = section(
            section_id="document",
            kind="document",
            label="Treść dokumentu",
            text=text,
        )
        sections = [item] if item else []

        return ExtractedContent(
            text=text,
            sections=sections,
            parser="libreoffice",
            status=(
                "ok"
                if len(text) >= 5
                else "empty_text"
            ),
            warnings=[],
            metrics={
                "section_count": len(
                    sections
                ),
            },
        )

    raise RuntimeError(
        "Brak programu antiword lub LibreOffice "
        "do odczytu formatu .doc."
    )


def extract_document(
    path: Path,
    extension: str,
) -> ExtractedContent:
    if extension == ".pdf":
        return extract_pdf(path)

    if extension == ".docx":
        return extract_docx(path)

    if extension == ".pptx":
        return extract_pptx(path)

    if extension == ".xlsx":
        return extract_xlsx(path)

    if extension == ".odt":
        return extract_odt(path)

    if extension == ".doc":
        return extract_legacy_doc(path)

    raise ValueError(
        "Brak parsera dla rozszerzenia: "
        f"{extension}"
    )


def deterministic_gzip_json(
    path: Path,
    payload: dict[str, Any],
) -> None:
    path.parent.mkdir(
        parents=True,
        exist_ok=True,
    )

    serialized = json.dumps(
        payload,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
    ).encode("utf-8")

    compressed = gzip.compress(
        serialized,
        compresslevel=9,
        mtime=0,
    )

    path.write_bytes(compressed)


def select_queue(
    queue: list[dict[str, Any]],
    *,
    sources: set[str],
    extensions: set[str],
    limit: int | None,
) -> list[dict[str, Any]]:
    selected: list[dict[str, Any]] = []

    for item in queue:
        source_id = str(
            item.get("source_id")
            or ""
        )
        url = normalize_url(
            item.get("effective_url")
            or item.get("final_url")
            or item.get("url")
        )
        extension = file_extension(
            url,
            item.get("filename"),
        )

        if (
            sources
            and source_id not in sources
        ):
            continue

        if (
            extensions
            and extension not in extensions
        ):
            continue

        selected.append(item)

    selected.sort(
        key=lambda item: (
            int(
                item.get("declared_filesize")
                or item.get("content_length")
                or 0
            ),
            str(
                item.get("effective_url")
                or item.get("url")
                or ""
            ),
        )
    )

    if limit is not None:
        selected = selected[:limit]

    return selected


def build_argument_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description=(
            "Pobiera i ekstrahuje tekst z dokumentów "
            "UEW, zachowując stan przyrostowy."
        )
    )

    parser.add_argument(
        "--limit",
        type=int,
        default=None,
        help=(
            "Maksymalna liczba dokumentów "
            "do przetworzenia."
        ),
    )
    parser.add_argument(
        "--source",
        action="append",
        default=[],
        help=(
            "Filtr source_id; można podać "
            "wielokrotnie."
        ),
    )
    parser.add_argument(
        "--extension",
        action="append",
        default=[],
        help=(
            "Filtr rozszerzenia, np. .pdf; "
            "można podać wielokrotnie."
        ),
    )
    parser.add_argument(
        "--max-file-size-mb",
        type=int,
        default=DEFAULT_MAX_FILE_SIZE_MB,
        help=(
            "Maksymalny rozmiar jednego pliku "
            "w MB."
        ),
    )
    parser.add_argument(
        "--timeout-seconds",
        type=int,
        default=DEFAULT_TIMEOUT_SECONDS,
        help="Limit czasu odczytu jednego pliku.",
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help=(
            "Ignoruj ETag i Last-Modified, "
            "pobierając dokument ponownie."
        ),
    )
    parser.add_argument(
        "--fail-on-error",
        action="store_true",
        help=(
            "Zakończ kodem 1, jeżeli "
            "wystąpi błąd ekstrakcji."
        ),
    )

    return parser


def main() -> None:
    arguments = build_argument_parser().parse_args()

    queue = load_jsonl(QUEUE_PATH)
    relationship_context = build_relationship_context()

    if not queue:
        raise FileNotFoundError(
            "Końcowa kolejka ekstrakcji jest "
            f"pusta lub nie istnieje: {QUEUE_PATH}"
        )

    previous_index = load_jsonl(INDEX_PATH)
    previous_state = load_jsonl(STATE_PATH)

    index_by_url = {
        normalize_url(
            item.get("effective_url")
            or item.get("url")
        ): item
        for item in previous_index
        if normalize_url(
            item.get("effective_url")
            or item.get("url")
        )
    }

    state_by_url = {
        normalize_url(
            item.get("effective_url")
            or item.get("url")
        ): item
        for item in previous_state
        if normalize_url(
            item.get("effective_url")
            or item.get("url")
        )
    }

    selected = select_queue(
        queue,
        sources={
            str(value).strip()
            for value in arguments.source
            if str(value).strip()
        },
        extensions={
            (
                str(value).strip().lower()
                if str(value).strip().startswith(".")
                else "." + str(value).strip().lower()
            )
            for value in arguments.extension
            if str(value).strip()
        },
        limit=arguments.limit,
    )

    full_run = (
        arguments.limit is None
        and not arguments.source
        and not arguments.extension
    )

    max_bytes = (
        arguments.max_file_size_mb
        * 1024
        * 1024
    )

    session = create_session()

    current_index_by_url = dict(index_by_url)
    current_state_by_url = dict(state_by_url)

    errors: list[dict[str, Any]] = []
    changes: list[dict[str, Any]] = []

    result_counts: Counter[str] = Counter()
    extraction_status_counts: Counter[str] = Counter()
    extension_counts: Counter[str] = Counter()
    source_counts: Counter[str] = Counter()

    downloaded_bytes_total = 0

    with tempfile.TemporaryDirectory() as directory:
        temporary_directory = Path(directory)

        for position, item in enumerate(
            selected,
            start=1,
        ):
            source_id = str(
                item.get("source_id")
                or "unknown"
            )
            requested_url = normalize_url(
                item.get("effective_url")
                or item.get("final_url")
                or item.get("url")
            )

            filename = safe_filename(
                item.get("filename"),
                fallback=(
                    Path(
                        unquote(
                            urlsplit(
                                requested_url
                            ).path
                        )
                    ).name
                    or (
                        f"attachment-"
                        f"{position:05d}"
                    )
                ),
            )
            extension = file_extension(
                requested_url,
                filename,
            )
            relationship_data = resolve_relationship_context(
                relationship_context,
                item,
                requested_url,
            )

            extension_counts[extension] += 1
            source_counts[source_id] += 1

            prior_state = state_by_url.get(
                requested_url
            )
            prior_index = index_by_url.get(
                requested_url
            )
            signature = queue_signature(item)

            temporary_path = (
                temporary_directory
                / (
                    f"{position:05d}-"
                    f"{filename}"
                )
            )

            print(
                f"[{position}/{len(selected)}] "
                f"{source_id}: {filename}"
            )

            try:
                download = download_document(
                    session=session,
                    url=requested_url,
                    destination=temporary_path,
                    previous_state=prior_state,
                    timeout_seconds=(
                        arguments.timeout_seconds
                    ),
                    max_bytes=max_bytes,
                    force=arguments.force,
                )

                if (
                    download["download_status"]
                    == "not_modified"
                ):
                    if (
                        prior_state is None
                        or prior_index is None
                    ):
                        raise RuntimeError(
                            "Serwer zwrócił 304, ale "
                            "brakuje poprzedniego stanu."
                        )

                    current_state_by_url[
                        requested_url
                    ] = {
                        **prior_state,
                        "checked_at": utc_now(),
                        "queue_signature": signature,
                        "http_status": 304,
                    }
                    current_index_by_url[
                        requested_url
                    ] = prior_index

                    result_counts["not_modified"] += 1
                    extraction_status_counts[
                        str(
                            prior_index.get(
                                "extraction_status"
                            )
                            or "unknown"
                        )
                    ] += 1

                    changes.append(
                        {
                            "effective_url": (
                                requested_url
                            ),
                            "change_type": (
                                "unchanged"
                            ),
                            "reason": (
                                "HTTP 304 Not Modified"
                            ),
                            "checked_at": utc_now(),
                        }
                    )
                    continue

                downloaded_bytes_total += int(
                    download.get(
                        "downloaded_bytes"
                    )
                    or 0
                )

                content_sha256 = sha256_file(
                    temporary_path
                )

                if (
                    prior_state
                    and prior_index
                    and prior_state.get(
                        "content_sha256"
                    ) == content_sha256
                    and prior_index.get(
                        "blob_path"
                    )
                    and Path(
                        str(
                            prior_index[
                                "blob_path"
                            ]
                        )
                    ).exists()
                ):
                    current_state_by_url[
                        requested_url
                    ] = {
                        **prior_state,
                        **download,
                        "effective_url": (
                            requested_url
                        ),
                        "queue_signature": signature,
                        "content_sha256": (
                            content_sha256
                        ),
                        "checked_at": utc_now(),
                    }
                    current_index_by_url[
                        requested_url
                    ] = {
                        **prior_index,
                        "checked_at": utc_now(),
                    }

                    result_counts[
                        "downloaded_unchanged"
                    ] += 1
                    extraction_status_counts[
                        str(
                            prior_index.get(
                                "extraction_status"
                            )
                            or "unknown"
                        )
                    ] += 1

                    changes.append(
                        {
                            "effective_url": (
                                requested_url
                            ),
                            "change_type": (
                                "unchanged"
                            ),
                            "reason": (
                                "SHA-256 bez zmian po "
                                "pobraniu kontrolnym"
                            ),
                            "content_sha256": (
                                content_sha256
                            ),
                            "checked_at": utc_now(),
                        }
                    )
                    continue

                extracted = extract_document(
                    temporary_path,
                    extension,
                )

                text_sha256 = sha256_bytes(
                    extracted.text.encode(
                        "utf-8"
                    )
                )

                blob_path = (
                    BLOB_DIRECTORY
                    / (
                        content_sha256
                        + ".json.gz"
                    )
                )

                blob_payload = {
                    "schema_version": (
                        SCHEMA_VERSION
                    ),
                    "attachment_id": (
                        item.get("candidate_id")
                        or item.get("queue_id")
                    ),
                    "queue_id": item.get(
                        "queue_id"
                    ),
                    "source_id": source_id,
                    "source_name": item.get(
                        "source_name"
                    ),
                    "source_priority": item.get(
                        "source_priority"
                    ),
                    "filename": filename,
                    "extension": extension,
                    "effective_url": (
                        requested_url
                    ),
                    "url_aliases": item.get(
                        "url_aliases",
                        [],
                    ),
                    "linked_from": relationship_data[
                        "linked_from"
                    ],
                    "linked_from_count": relationship_data[
                        "linked_from_count"
                    ],
                    "attachment_metadata": relationship_data[
                        "attachment_metadata"
                    ],
                    "content_type": (
                        download.get(
                            "content_type"
                        )
                        or item.get(
                            "content_type"
                        )
                        or item.get(
                            "declared_mime_type"
                        )
                    ),
                    "content_sha256": (
                        content_sha256
                    ),
                    "text_sha256": (
                        text_sha256
                    ),
                    "extracted_at": utc_now(),
                    "parser": extracted.parser,
                    "extraction_status": (
                        extracted.status
                    ),
                    "warnings": (
                        extracted.warnings
                    ),
                    "metrics": (
                        extracted.metrics
                    ),
                    "char_count": len(
                        extracted.text
                    ),
                    "word_count": word_count(
                        extracted.text
                    ),
                    "section_count": len(
                        extracted.sections
                    ),
                    "text": extracted.text,
                    "sections": (
                        extracted.sections
                    ),
                }

                deterministic_gzip_json(
                    blob_path,
                    blob_payload,
                )

                index_record = {
                    "schema_version": (
                        SCHEMA_VERSION
                    ),
                    "attachment_id": (
                        item.get("candidate_id")
                        or item.get("queue_id")
                    ),
                    "queue_id": item.get(
                        "queue_id"
                    ),
                    "source_id": source_id,
                    "source_name": item.get(
                        "source_name"
                    ),
                    "source_priority": item.get(
                        "source_priority"
                    ),
                    "filename": filename,
                    "extension": extension,
                    "effective_url": (
                        requested_url
                    ),
                    "url_aliases": item.get(
                        "url_aliases",
                        [],
                    ),
                    "linked_from": relationship_data[
                        "linked_from"
                    ],
                    "linked_from_count": relationship_data[
                        "linked_from_count"
                    ],
                    "attachment_metadata": relationship_data[
                        "attachment_metadata"
                    ],
                    "content_sha256": (
                        content_sha256
                    ),
                    "text_sha256": (
                        text_sha256
                    ),
                    "blob_path": str(
                        blob_path
                    ),
                    "parser": (
                        extracted.parser
                    ),
                    "extraction_status": (
                        extracted.status
                    ),
                    "char_count": len(
                        extracted.text
                    ),
                    "word_count": word_count(
                        extracted.text
                    ),
                    "section_count": len(
                        extracted.sections
                    ),
                    "metrics": (
                        extracted.metrics
                    ),
                    "warnings": (
                        extracted.warnings
                    ),
                    "extracted_at": (
                        blob_payload[
                            "extracted_at"
                        ]
                    ),
                    "checked_at": utc_now(),
                }

                state_record = {
                    "schema_version": (
                        SCHEMA_VERSION
                    ),
                    "effective_url": (
                        requested_url
                    ),
                    "queue_signature": (
                        signature
                    ),
                    "filename": filename,
                    "extension": extension,
                    "source_id": source_id,
                    "content_sha256": (
                        content_sha256
                    ),
                    "text_sha256": (
                        text_sha256
                    ),
                    "blob_path": str(
                        blob_path
                    ),
                    "etag": download.get(
                        "etag"
                    ),
                    "last_modified": (
                        download.get(
                            "last_modified"
                        )
                    ),
                    "content_length": (
                        download.get(
                            "content_length"
                        )
                    ),
                    "content_type": (
                        download.get(
                            "content_type"
                        )
                    ),
                    "http_status": (
                        download.get(
                            "http_status"
                        )
                    ),
                    "final_url": (
                        download.get(
                            "final_url"
                        )
                    ),
                    "redirect_count": (
                        download.get(
                            "redirect_count"
                        )
                    ),
                    "extraction_status": (
                        extracted.status
                    ),
                    "checked_at": utc_now(),
                    "extracted_at": (
                        blob_payload[
                            "extracted_at"
                        ]
                    ),
                }

                current_index_by_url[
                    requested_url
                ] = index_record
                current_state_by_url[
                    requested_url
                ] = state_record

                change_type = (
                    "new"
                    if prior_state is None
                    else "changed"
                )

                changes.append(
                    {
                        "effective_url": (
                            requested_url
                        ),
                        "change_type": (
                            change_type
                        ),
                        "content_sha256": (
                            content_sha256
                        ),
                        "text_sha256": (
                            text_sha256
                        ),
                        "blob_path": str(
                            blob_path
                        ),
                        "extraction_status": (
                            extracted.status
                        ),
                        "checked_at": utc_now(),
                    }
                )

                result_counts[
                    change_type
                ] += 1
                extraction_status_counts[
                    extracted.status
                ] += 1

            except Exception as error:
                error_record = {
                    "effective_url": (
                        requested_url
                    ),
                    "source_id": source_id,
                    "filename": filename,
                    "extension": extension,
                    "error": str(error),
                    "error_type": (
                        type(error).__name__
                    ),
                    "failed_at": utc_now(),
                }

                errors.append(error_record)
                changes.append(
                    {
                        "effective_url": (
                            requested_url
                        ),
                        "change_type": "failed",
                        "error": str(error),
                        "checked_at": utc_now(),
                    }
                )
                result_counts["failed"] += 1

                print(
                    f"  BŁĄD: {error}"
                )

    session.close()

    queue_urls = {
        normalize_url(
            item.get("effective_url")
            or item.get("final_url")
            or item.get("url")
        )
        for item in queue
        if normalize_url(
            item.get("effective_url")
            or item.get("final_url")
            or item.get("url")
        )
    }

    deleted_count = 0

    if full_run:
        for url in sorted(
            set(current_index_by_url)
            - queue_urls
        ):
            current_index_by_url.pop(
                url,
                None,
            )
            current_state_by_url.pop(
                url,
                None,
            )
            deleted_count += 1
            changes.append(
                {
                    "effective_url": url,
                    "change_type": "deleted",
                    "reason": (
                        "Adres nie występuje już "
                        "w końcowej kolejce."
                    ),
                    "checked_at": utc_now(),
                }
            )

    current_index = sorted(
        current_index_by_url.values(),
        key=lambda item: str(
            item.get("effective_url")
            or ""
        ),
    )
    current_state = sorted(
        current_state_by_url.values(),
        key=lambda item: str(
            item.get("effective_url")
            or ""
        ),
    )

    write_jsonl(
        INDEX_PATH,
        current_index,
    )
    write_jsonl(
        STATE_PATH,
        current_state,
    )
    write_jsonl(
        ERRORS_PATH,
        errors,
    )
    write_jsonl(
        CHANGES_PATH,
        changes,
    )

    status = {
        "schema_version": SCHEMA_VERSION,
        "generated_at": utc_now(),
        "queue_count": len(queue),
        "selected_count": len(selected),
        "full_run": full_run,
        "previous_index_count": len(
            previous_index
        ),
        "current_index_count": len(
            current_index
        ),
        "processed_successfully_count": (
            len(selected)
            - result_counts["failed"]
        ),
        "error_count": len(errors),
        "deleted_count": deleted_count,
        "result_counts": dict(
            sorted(result_counts.items())
        ),
        "extraction_status_counts": dict(
            sorted(
                extraction_status_counts.items()
            )
        ),
        "selected_extension_counts": dict(
            sorted(extension_counts.items())
        ),
        "selected_source_counts": dict(
            sorted(source_counts.items())
        ),
        "downloaded_bytes": (
            downloaded_bytes_total
        ),
        "downloaded_megabytes": round(
            downloaded_bytes_total
            / (1024 * 1024),
            2,
        ),
        "files": {
            "index": str(INDEX_PATH),
            "state": str(STATE_PATH),
            "errors": str(ERRORS_PATH),
            "changes": str(CHANGES_PATH),
            "blob_directory": str(
                BLOB_DIRECTORY
            ),
        },
    }

    STATUS_PATH.parent.mkdir(
        parents=True,
        exist_ok=True,
    )
    STATUS_PATH.write_text(
        json.dumps(
            status,
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    print("Ekstrakcja zakończona.")
    print(
        "Wybrane dokumenty: "
        f"{len(selected)}"
    )
    print(
        "Błędy: "
        f"{len(errors)}"
    )
    print(
        "Aktualny indeks tekstów: "
        f"{len(current_index)}"
    )
    print(
        "Pobrano: "
        f"{status['downloaded_megabytes']} MB"
    )

    if arguments.fail_on_error and errors:
        raise SystemExit(1)


if __name__ == "__main__":
    main()
